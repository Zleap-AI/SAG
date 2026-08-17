"""上传文件 → Markdown 的路由、缓存与 302 MinerU 适配。"""

from __future__ import annotations

import asyncio
import io
import logging
import zipfile
from pathlib import Path
from types import SimpleNamespace
from typing import Any

import httpx
import pytest

from sag_api.core.config import Settings
from sag_api.core.errors import (
    ConfigurationError,
    ServiceUnavailableError,
    UpstreamError,
)
from sag_api.parsing import service
from sag_api.parsing.mineru import MinerUClient, _assert_public_host, _interpret_poll_payload
from sag_api.parsing.service import PreparedDocument
from sag_api.sag.dto import ProcessCheckpoint, ProcessOutcome


def _settings(**overrides: Any) -> Settings:
    return Settings(
        _env_file=None,
        data_dir="/tmp/sag-test-engine",
        upload_dir="/tmp/sag-test-uploads",
        **overrides,
    )


@pytest.mark.asyncio
async def test_parser_routes_markdown_and_markitdown_with_cache(tmp_path, monkeypatch):
    markdown = tmp_path / "already.md"
    markdown.write_text("# Already\n", encoding="utf-8")
    direct = await service.prepare_document(str(markdown), _settings())
    assert direct.path == str(markdown)
    assert direct.provider == "original"

    source = tmp_path / "notes.docx"
    source.write_bytes(b"fake-office")
    calls: list[str] = []

    def convert(path: str) -> str:
        calls.append(path)
        return "# Converted\n\nhello"

    monkeypatch.setattr(service, "_markitdown_sync", convert)
    first = await service.prepare_document(str(source), _settings())
    second = await service.prepare_document(str(source), _settings())

    assert first.provider == "markitdown" and first.path.endswith(".parsed.markitdown.md")
    assert Path(first.path).read_text(encoding="utf-8") == "# Converted\n\nhello\n"
    assert second.cached is True and second.path == first.path
    assert calls == [str(source)]


@pytest.mark.asyncio
@pytest.mark.parametrize(
    ("settings_overrides", "resume_state", "expected_metadata", "expected_public_model"),
    [
        (
            {
                "mineru_provider": "official",
                "mineru_official_model": "pipeline",
                "mineru_api_key": "official-key",
            },
            None,
            {"mineru_service": "official", "mineru_model": "pipeline"},
            "pipeline",
        ),
        (
            {
                "mineru_provider": "302",
                "mineru_version": "2.0",
                "mineru_api_key": "302-key",
            },
            {"provider": "mineru", "status": "done"},
            {"mineru_service": "302", "mineru_version": "2.0"},
            "pipeline",
        ),
    ],
)
async def test_mineru_cache_hit_emits_authoritative_metadata(
    tmp_path, settings_overrides, resume_state, expected_metadata, expected_public_model
):
    source = tmp_path / "cached.pdf"
    source.write_bytes(b"%PDF-fake")
    settings = _settings(
        document_parser="mineru",
        mineru_base_url="https://mineru.example.test",
        **settings_overrides,
    )
    signature = service._signature("mineru", settings)
    cache_path = Path(f"{source}.parsed.{signature}.md")
    cache_path.write_text("# Cached result\n", encoding="utf-8")
    states: list[dict[str, Any]] = []

    prepared = await service.prepare_document(
        str(source),
        settings,
        state=resume_state,
        on_state=lambda state: _append_state(states, state),
    )

    assert prepared.cached is True
    assert prepared.provider == "mineru"
    assert states[-1]["provider"] == "mineru"
    assert states[-1]["status"] == "done"
    assert states[-1]["cache_path"] == str(cache_path)
    assert {key: states[-1][key] for key in expected_metadata} == expected_metadata

    from sag_api.jobs import tasks

    public_values = tasks._prepared_parser_values(prepared, states[-1])
    assert public_values["mineru_provider"] == expected_metadata["mineru_service"]
    assert public_values["mineru_model"] == expected_public_model


async def _append_state(states: list[dict[str, Any]], state: dict[str, Any]) -> None:
    states.append(state)


@pytest.mark.asyncio
async def test_legacy_gb18030_text_is_normalized_without_markitdown(tmp_path, monkeypatch):
    source = tmp_path / "骆驼祥子.txt"
    expected = "《骆驼祥子》\r\n作者：老舍\r\n正文只有一个损坏字节："
    source.write_bytes(expected.encode("gb18030") + b"\xff")
    stale_cache = Path(f"{source}.parsed.markitdown.md")
    stale_cache.write_text("None\n", encoding="utf-8")

    def should_not_run(_path: str) -> str:
        raise AssertionError("plain text should use Muse's text decoder")

    monkeypatch.setattr(service, "_markitdown_sync", should_not_run)
    parsed = await service.prepare_document(str(source), _settings())
    markdown = Path(parsed.path).read_text(encoding="utf-8")

    assert parsed.cached is False
    assert parsed.provider == "markitdown"
    assert markdown.startswith(expected.replace("\r\n", "\n"))
    assert markdown.count("�") == 1
    assert markdown != "None\n"


@pytest.mark.asyncio
async def test_markitdown_none_sentinel_is_rejected(tmp_path, monkeypatch):
    source = tmp_path / "broken.docx"
    source.write_bytes(b"fake-office")
    monkeypatch.setattr(service, "_markitdown_sync", lambda _path: "None")

    with pytest.raises(Exception, match="未从文件中解析出有效文本"):
        await service.prepare_document(str(source), _settings())

    assert not Path(f"{source}.parsed.markitdown.md").exists()


@pytest.mark.asyncio
async def test_only_pdf_uses_configured_mineru(tmp_path, monkeypatch):
    settings = _settings(
        document_parser="auto",
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )
    seen: list[str] = []

    class FakeMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None, should_pause=None):
            seen.append(path)
            return "# From MinerU\n"

    monkeypatch.setattr(service, "MinerUClient", FakeMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", lambda path: "# From MarkItDown\n")

    pdf = tmp_path / "paper.pdf"
    pdf.write_bytes(b"%PDF-fake")
    docx = tmp_path / "paper.docx"
    docx.write_bytes(b"fake-office")

    parsed_pdf = await service.prepare_document(str(pdf), settings)
    parsed_docx = await service.prepare_document(str(docx), settings)
    assert parsed_pdf.provider == "mineru"
    assert parsed_docx.provider == "markitdown"
    assert seen == [str(pdf)]


@pytest.mark.asyncio
async def test_pdf_without_complete_mineru_config_falls_back_to_markitdown(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    monkeypatch.setattr(service, "_markitdown_sync", lambda path: "# Local PDF\n")

    parsed = await service.prepare_document(
        str(source),
        _settings(document_parser="mineru", mineru_base_url="https://api.302.ai"),
    )
    assert parsed.provider == "markitdown"


@pytest.mark.asyncio
async def test_mineru_failure_falls_back_to_markitdown_and_reuses_cache(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    mineru_calls = 0
    markitdown_calls = 0

    class FailingMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None, should_pause=None):
            nonlocal mineru_calls
            mineru_calls += 1
            if on_state:
                await on_state({**(state or {}), "task_id": "task-unavailable"})
            raise ServiceUnavailableError("No available models currently")

    def convert(_path: str) -> str:
        nonlocal markitdown_calls
        markitdown_calls += 1
        return "# Local fallback\n"

    monkeypatch.setattr(service, "MinerUClient", FailingMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", convert)
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )
    states: list[dict[str, Any]] = []

    first = await service.prepare_document(str(source), settings, on_state=lambda state: _record(states, state))
    second = await service.prepare_document(
        str(source),
        settings,
        state=states[-1],
        on_state=lambda state: _record(states, state),
    )

    assert first.provider == "markitdown" and first.path.endswith(".parsed.markitdown.md")
    assert Path(first.path).read_text(encoding="utf-8") == "# Local fallback\n"
    assert second.path == first.path and second.cached is True
    assert mineru_calls == 1 and markitdown_calls == 1
    assert states[-1]["provider"] == "mineru"
    assert states[-1]["status"] == "fallback_done"
    assert states[-1]["fallback"]["provider"] == "markitdown"
    assert states[-1]["fallback"]["mineru_error"] == "No available models currently"
    assert states[-1]["fallback"]["status"] == "done"
    assert first.fallback_from == "mineru"
    assert first.fallback_error == "No available models currently"


@pytest.mark.asyncio
async def test_document_fails_only_when_mineru_and_markitdown_both_fail(tmp_path, monkeypatch):
    source = tmp_path / "broken.pdf"
    source.write_bytes(b"%PDF-broken")

    class FailingMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None, should_pause=None):
            raise ServiceUnavailableError("remote parser failed")

    def fail_locally(_path: str) -> str:
        raise RuntimeError("local parser failed")

    monkeypatch.setattr(service, "MinerUClient", FailingMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", fail_locally)
    states: list[dict[str, Any]] = []
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )

    with pytest.raises(ServiceUnavailableError, match="MinerU.*MarkItDown"):
        await service.prepare_document(str(source), settings, on_state=lambda state: _record(states, state))
    assert states[-1]["status"] == "fallback_failed"
    assert states[-1]["fallback"]["markitdown_error"].endswith("local parser failed")


@pytest.mark.asyncio
async def test_mineru_state_callback_failure_does_not_trigger_markitdown(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    markitdown_calls = 0

    class MinerUWithState:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None, should_pause=None):
            assert on_state is not None
            await on_state({**(state or {}), "task_id": "task-1"})
            return "# Never reached\n"

    def convert(_path: str) -> str:
        nonlocal markitdown_calls
        markitdown_calls += 1
        return "# Should not run\n"

    state_writes = 0

    async def fail_to_persist(_state: dict[str, Any]) -> None:
        nonlocal state_writes
        state_writes += 1
        if state_writes == 2:
            raise RuntimeError("database commit failed")

    monkeypatch.setattr(service, "MinerUClient", MinerUWithState)
    monkeypatch.setattr(service, "_markitdown_sync", convert)
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )

    with pytest.raises(RuntimeError, match="database commit failed"):
        await service.prepare_document(str(source), settings, on_state=fail_to_persist)
    assert state_writes == 2
    assert markitdown_calls == 0


@pytest.mark.asyncio
async def test_changed_mineru_config_retries_remote_after_cached_fallback(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    mineru_calls = 0

    class FailingMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None, should_pause=None):
            nonlocal mineru_calls
            mineru_calls += 1
            raise UpstreamError("remote unavailable")

    monkeypatch.setattr(service, "MinerUClient", FailingMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", lambda path: "# Local fallback\n")
    states: list[dict[str, Any]] = []
    first_settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-first",
    )
    second_settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-changed",
    )

    await service.prepare_document(str(source), first_settings, on_state=lambda state: _record(states, state))
    previous_state = states[-1]
    await service.prepare_document(
        str(source),
        second_settings,
        state=previous_state,
        on_state=lambda state: _record(states, state),
    )

    assert mineru_calls == 2
    assert states[-1]["key_fingerprint"] != previous_state["key_fingerprint"]


@pytest.mark.asyncio
async def test_concurrent_pdf_parsing_creates_only_one_mineru_task(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    started = asyncio.Event()
    release = asyncio.Event()
    calls = 0

    class FakeMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None, should_pause=None):
            nonlocal calls
            calls += 1
            started.set()
            await release.wait()
            return "# Parsed once\n"

    monkeypatch.setattr(service, "MinerUClient", FakeMinerU)
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )
    first = asyncio.create_task(service.prepare_document(str(source), settings))
    second = asyncio.create_task(service.prepare_document(str(source), settings))
    await started.wait()
    release.set()
    results = await asyncio.gather(first, second)

    assert calls == 1
    assert results[0].path == results[1].path
    assert sorted(result.cached for result in results) == [False, True]


@pytest.mark.asyncio
async def test_concurrent_mineru_failure_creates_one_task_and_one_fallback(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    mineru_calls = 0
    markitdown_calls = 0

    class FailingMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None, should_pause=None):
            nonlocal mineru_calls
            mineru_calls += 1
            await asyncio.sleep(0.01)
            raise ServiceUnavailableError("No available models currently")

    def convert(_path: str) -> str:
        nonlocal markitdown_calls
        markitdown_calls += 1
        return "# One local fallback\n"

    monkeypatch.setattr(service, "MinerUClient", FailingMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", convert)
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )

    results = await asyncio.gather(
        service.prepare_document(str(source), settings),
        service.prepare_document(str(source), settings),
    )

    assert mineru_calls == 1 and markitdown_calls == 1
    assert results[0].path == results[1].path
    assert all(result.provider == "markitdown" for result in results)
    assert sorted(result.cached for result in results) == [False, True]


@pytest.mark.asyncio
async def test_document_job_sends_parsed_markdown_to_engine(monkeypatch):
    from sag_api.db.models import Document, Source
    from sag_api.enums import DocumentStatus
    from sag_api.jobs import tasks

    document = SimpleNamespace(
        id="doc-1",
        source_id="source-1",
        filename="original.pdf",
        storage_path="/uploads/original.pdf",
        status=None,
        error="previous attempt failed",
        chunk_count=0,
        event_count=0,
        progress=0,
        token_usage=0,
        sag_source_id=None,
    )
    source = SimpleNamespace(
        id="source-1",
        sag_source_config_id="sag-source-1",
        chunk_count=0,
        event_count=0,
    )
    job = SimpleNamespace(id="job-1", document_id="doc-1", progress=0.0, payload={})

    class FakeSession:
        async def get(self, model, _id):
            return document if model is Document else source if model is Source else None

        async def commit(self):
            pass

        async def execute(self, _statement):
            document.status = DocumentStatus.READY
            document.chunk_count = 2
            document.event_count = 1
            document.sag_source_id = "engine-doc"
            document.progress = 100
            document.token_usage = 2468
            return SimpleNamespace(rowcount=1)

        async def scalar(self, _statement):
            return None

        async def refresh(self, _instance, attribute_names=None):
            pass

    prepared_calls: list[str] = []
    stage_errors: list[tuple[str, str | None]] = []

    async def fake_prepare(path, settings, *, state=None, on_state=None, should_pause=None):
        prepared_calls.append(path)
        return PreparedDocument("/uploads/original.pdf.parsed.markitdown.md", "markitdown")

    class FakeEngineManager:
        seen_path = ""

        async def process_document(
            self,
            source_config_id,
            path,
            *,
            source,
            on_stage,
            checkpoint,
            on_checkpoint,
            should_pause,
            max_concurrency,
            document_title,
        ):
            self.seen_path = path
            assert max_concurrency == tasks.settings.document_extract_concurrency
            assert document_title == "original"
            await on_stage("loading")
            stage_errors.append(("loading", document.error))
            await on_checkpoint(
                ProcessCheckpoint(
                    source_id="engine-doc",
                    chunk_ids=["chunk-1", "chunk-2"],
                    processed_chunk_ids=["chunk-1"],
                    event_count=1,
                    event_ids=["event-1"],
                    token_usage=1234,
                )
            )
            await on_stage("extracting")
            stage_errors.append(("extracting", document.error))
            return ProcessOutcome(
                source_id="engine-doc",
                chunk_count=2,
                event_count=1,
                chunk_ids=["chunk-1", "chunk-2"],
                processed_chunk_ids=["chunk-1", "chunk-2"],
                token_usage=2468,
            )

    monkeypatch.setattr(tasks, "prepare_document", fake_prepare)
    engine = FakeEngineManager()
    await tasks._process_document_unlocked(FakeSession(), job, engine_manager=engine)

    assert prepared_calls == ["/uploads/original.pdf"]
    assert stage_errors == [("loading", None), ("extracting", None)]
    assert engine.seen_path.endswith(".md")
    assert document.status.value == "ready"
    assert document.chunk_count == 2 and document.event_count == 1
    assert document.progress == 100 and document.token_usage == 2468
    assert document.parser_provider == "markitdown"
    assert document.parser_status == "done"


@pytest.mark.asyncio
@pytest.mark.parametrize(
    ("parser_state", "expected_provider", "expected_model"),
    [
        (
            {
                "provider": "mineru",
                "mineru_service": "official",
                "mineru_model": "pipeline",
            },
            "official",
            "pipeline",
        ),
        (
            {
                "provider": "mineru",
                "mineru_service": "302",
                "mineru_version": "2.5",
            },
            "302",
            "2.5",
        ),
        (
            {
                "provider": "mineru",
                "mineru_service": "302",
                "mineru_version": "2.0",
                "base_url": "https://mineru-gateway.example.test",
            },
            "302",
            "pipeline",
        ),
    ],
)
async def test_document_job_persists_successful_mineru_outcome(
    monkeypatch, parser_state, expected_provider, expected_model
):
    from sag_api.db.models import Document, Source
    from sag_api.enums import DocumentStatus
    from sag_api.jobs import tasks

    document = SimpleNamespace(
        id="doc-mineru",
        source_id="source-1",
        filename="paper.pdf",
        storage_path="/uploads/paper.pdf",
        status=DocumentStatus.PENDING,
        error=None,
        chunk_count=0,
        event_count=0,
        progress=0,
        token_usage=0,
        sag_source_id=None,
    )
    source = SimpleNamespace(id="source-1", sag_source_config_id="sag-source-1")
    job = SimpleNamespace(id="job-mineru", document_id=document.id, progress=0.0, payload={})

    class FakeSession:
        async def get(self, model, _id):
            return document if model is Document else source if model is Source else None

        async def refresh(self, _instance, attribute_names=None):
            pass

        async def commit(self):
            pass

        async def execute(self, _statement):
            return SimpleNamespace(rowcount=1)

    async def fake_prepare(path, settings, *, state=None, on_state=None, should_pause=None):
        assert on_state is not None
        await on_state(
            {
                **parser_state,
                "task_id": "task-1",
                "status": "running",
            }
        )
        return PreparedDocument(f"{path}.parsed.mineru.md", "mineru")

    class FakeEngineManager:
        async def process_document(self, *args, **kwargs):
            return ProcessOutcome(
                source_id="engine-doc",
                chunk_count=1,
                event_count=1,
                chunk_ids=["chunk-1"],
                processed_chunk_ids=["chunk-1"],
                token_usage=123,
            )

    async def no_op_touch(*args):
        return None

    monkeypatch.setattr(tasks, "prepare_document", fake_prepare)
    monkeypatch.setattr(tasks, "touch_source_revision", no_op_touch)

    await tasks._process_document_unlocked(FakeSession(), job, engine_manager=FakeEngineManager())

    assert document.parser_provider == "mineru"
    assert document.mineru_provider == expected_provider
    assert document.mineru_model == expected_model
    assert document.parser_status == "done"
    assert document.fallback_from is None
    assert document.fallback_reason is None
    assert job.payload["document_parser"]["task_id"] == "task-1"


@pytest.mark.asyncio
async def test_document_job_persists_mineru_markitdown_fallback(monkeypatch, caplog):
    from sag_api.db.models import Document, Source
    from sag_api.enums import DocumentStatus
    from sag_api.jobs import tasks

    document = SimpleNamespace(
        id="doc-fallback",
        source_id="source-1",
        filename="paper.pdf",
        storage_path="/uploads/paper.pdf",
        status=DocumentStatus.PENDING,
        error=None,
        chunk_count=0,
        event_count=0,
        progress=0,
        token_usage=0,
        sag_source_id=None,
    )
    source = SimpleNamespace(id="source-1", sag_source_config_id="sag-source-1")
    job = SimpleNamespace(id="job-fallback", document_id=document.id, progress=0.0, payload={})
    raw_reason = "MinerU failed with sk-secret123 at https://files.example/result?token=signed-value"

    class FakeSession:
        async def get(self, model, _id):
            return document if model is Document else source if model is Source else None

        async def refresh(self, _instance, attribute_names=None):
            pass

        async def commit(self):
            pass

        async def execute(self, _statement):
            return SimpleNamespace(rowcount=1)

    async def fake_prepare(path, settings, *, state=None, on_state=None, should_pause=None):
        assert on_state is not None
        await on_state(
            {
                "provider": "mineru",
                "mineru_service": "302",
                "mineru_version": "2.5",
                "base_url": "https://api.302ai.cn",
                "status": "fallback_done",
                "fallback": {
                    "provider": "markitdown",
                    "status": "done",
                    "mineru_error": raw_reason,
                },
            }
        )
        return PreparedDocument(
            f"{path}.parsed.markitdown.md",
            "markitdown",
            fallback_from="mineru",
            fallback_error=raw_reason,
        )

    class FakeEngineManager:
        async def process_document(self, *args, **kwargs):
            return ProcessOutcome(
                source_id="engine-doc",
                chunk_count=1,
                event_count=0,
                chunk_ids=["chunk-1"],
                processed_chunk_ids=["chunk-1"],
                token_usage=0,
            )

    async def no_op_touch(*args):
        return None

    monkeypatch.setattr(tasks, "prepare_document", fake_prepare)
    monkeypatch.setattr(tasks, "touch_source_revision", no_op_touch)

    caplog.set_level(logging.WARNING, logger="sag.jobs")
    await tasks._process_document_unlocked(
        FakeSession(), job, engine_manager=FakeEngineManager()
    )

    assert document.parser_provider == "markitdown"
    assert document.mineru_provider == "302"
    assert document.mineru_model == "2.5"
    assert document.parser_status == "fallback"
    assert document.fallback_from == "mineru"
    assert "sk-secret123" not in document.fallback_reason
    assert "signed-value" not in document.fallback_reason
    warning_text = " ".join(record.getMessage() for record in caplog.records)
    assert "sk-secret123" not in warning_text
    assert "signed-value" not in warning_text
    assert "[REDACTED]" in warning_text
    assert job.payload["document_parser"]["fallback"]["mineru_error"] == raw_reason


@pytest.mark.asyncio
async def test_document_job_redacts_parser_failure_from_public_error(
    monkeypatch, caplog
):
    from sag_api.db.models import Document, Source
    from sag_api.enums import DocumentStatus
    from sag_api.jobs import tasks

    raw_reason = (
        "Parser failed with Bearer secret-token and sk-secret123 at "
        "https://files.example/result?token=signed-value " + "detail " * 100
    )
    document = SimpleNamespace(
        id="doc-parser-failed",
        source_id="source-1",
        filename="paper.pdf",
        storage_path="/uploads/paper.pdf",
        status=DocumentStatus.PENDING,
        error=None,
        chunk_count=0,
        event_count=0,
        progress=0,
        token_usage=0,
        sag_source_id=None,
    )
    source = SimpleNamespace(id="source-1", sag_source_config_id="sag-source-1")
    job = SimpleNamespace(
        id="job-parser-failed",
        document_id=document.id,
        progress=0.0,
        payload={},
    )

    class FakeSession:
        async def get(self, model, _id):
            return document if model is Document else source if model is Source else None

        async def refresh(self, _instance, attribute_names=None):
            pass

        async def commit(self):
            pass

        async def execute(self, statement):
            values = {
                column.key: bound.value for column, bound in statement._values.items()
            }
            if "error" in values:
                document.error = values["error"]
            return SimpleNamespace(rowcount=1)

    async def fake_prepare(path, settings, *, state=None, on_state=None, should_pause=None):
        assert on_state is not None
        await on_state(
            {
                "provider": "mineru",
                "mineru_service": "302",
                "mineru_version": "2.5",
                "status": "fallback_failed",
                "fallback": {
                    "provider": "markitdown",
                    "status": "failed",
                    "mineru_error": raw_reason,
                    "markitdown_error": raw_reason,
                },
            }
        )
        raise ServiceUnavailableError(raw_reason)

    class FakeEngineManager:
        async def process_document(self, *args, **kwargs):
            raise AssertionError("engine must not run after parser failure")

    monkeypatch.setattr(tasks, "prepare_document", fake_prepare)
    caplog.set_level(logging.WARNING, logger="sag.jobs")

    with pytest.raises(ServiceUnavailableError):
        await tasks._process_document_unlocked(
            FakeSession(), job, engine_manager=FakeEngineManager()
        )

    assert document.parser_status == "failed"
    assert document.fallback_reason == document.error
    assert document.error is not None and len(document.error) <= 300
    assert "secret-token" not in document.error
    assert "sk-secret123" not in document.error
    assert "signed-value" not in document.error
    warning_text = " ".join(record.getMessage() for record in caplog.records)
    assert "secret-token" not in warning_text
    assert "sk-secret123" not in warning_text
    assert "signed-value" not in warning_text


@pytest.mark.asyncio
async def test_document_job_preserves_engine_error_before_first_checkpoint(
    monkeypatch, caplog
):
    from sag_api.db.models import Document, Source
    from sag_api.enums import DocumentStatus
    from sag_api.jobs import tasks

    raw_reason = "engine extraction failed before checkpoint " + "diagnostic " * 40
    document = SimpleNamespace(
        id="doc-engine-failed",
        source_id="source-1",
        filename="paper.pdf",
        storage_path="/uploads/paper.pdf",
        status=DocumentStatus.PENDING,
        error=None,
        chunk_count=0,
        event_count=0,
        progress=0,
        token_usage=0,
        sag_source_id=None,
    )
    source = SimpleNamespace(id="source-1", sag_source_config_id="sag-source-1")
    job = SimpleNamespace(
        id="job-engine-failed",
        document_id=document.id,
        progress=0.0,
        payload={},
    )

    class FakeSession:
        async def get(self, model, _id):
            return document if model is Document else source if model is Source else None

        async def refresh(self, _instance, attribute_names=None):
            pass

        async def commit(self):
            pass

        async def execute(self, statement):
            values = {
                column.key: bound.value for column, bound in statement._values.items()
            }
            if "error" in values:
                document.error = values["error"]
            return SimpleNamespace(rowcount=1)

    async def fake_prepare(path, settings, *, state=None, on_state=None, should_pause=None):
        assert on_state is not None
        await on_state(
            {
                "provider": "mineru",
                "mineru_service": "official",
                "mineru_model": "pipeline",
                "status": "done",
            }
        )
        return PreparedDocument(f"{path}.parsed.mineru.md", "mineru")

    class FakeEngineManager:
        async def process_document(self, *args, **kwargs):
            assert not kwargs["checkpoint"].chunk_ids
            raise RuntimeError(raw_reason)

    monkeypatch.setattr(tasks, "prepare_document", fake_prepare)
    caplog.set_level(logging.WARNING, logger="sag.jobs")

    with pytest.raises(RuntimeError, match="engine extraction failed"):
        await tasks._process_document_unlocked(
            FakeSession(), job, engine_manager=FakeEngineManager()
        )

    assert document.error == raw_reason
    assert len(document.error) > 300
    warning_text = " ".join(record.getMessage() for record in caplog.records)
    assert raw_reason in warning_text


@pytest.mark.asyncio
async def test_document_job_preserves_engine_error_on_resumed_checkpoint(
    monkeypatch, caplog
):
    from sag_api.db.models import Document, Source
    from sag_api.enums import DocumentStatus
    from sag_api.jobs import tasks

    raw_reason = "resumed engine extraction failed " + "checkpoint diagnostic " * 30
    document = SimpleNamespace(
        id="doc-resumed-engine-failed",
        source_id="source-1",
        filename="paper.pdf",
        storage_path="/uploads/paper.pdf",
        status=DocumentStatus.EXTRACTING,
        error=None,
        chunk_count=1,
        event_count=0,
        progress=20,
        token_usage=0,
        sag_source_id=None,
    )
    source = SimpleNamespace(id="source-1", sag_source_config_id="sag-source-1")
    job = SimpleNamespace(
        id="job-resumed-engine-failed",
        document_id=document.id,
        progress=0.2,
        payload={
            "process_checkpoint": {
                "chunk_ids": ["chunk-1"],
                "processed_chunk_ids": [],
            }
        },
    )

    class FakeSession:
        async def get(self, model, _id):
            return document if model is Document else source if model is Source else None

        async def refresh(self, _instance, attribute_names=None):
            pass

        async def commit(self):
            pass

        async def execute(self, statement):
            values = {
                column.key: bound.value for column, bound in statement._values.items()
            }
            if "error" in values:
                document.error = values["error"]
            return SimpleNamespace(rowcount=1)

    async def unexpected_prepare(*args, **kwargs):
        pytest.fail("resumed jobs must skip parser preparation")

    class FakeEngineManager:
        async def process_document(self, *args, **kwargs):
            assert kwargs["checkpoint"].chunk_ids == ["chunk-1"]
            raise RuntimeError(raw_reason)

    monkeypatch.setattr(tasks, "prepare_document", unexpected_prepare)
    caplog.set_level(logging.WARNING, logger="sag.jobs")

    with pytest.raises(RuntimeError, match="resumed engine extraction failed"):
        await tasks._process_document_unlocked(
            FakeSession(), job, engine_manager=FakeEngineManager()
        )

    assert document.error == raw_reason
    assert len(document.error) > 300
    warning_text = " ".join(record.getMessage() for record in caplog.records)
    assert raw_reason in warning_text


class _FakeAsyncClient:
    responses: list[httpx.Response] = []
    calls: list[tuple[str, str, dict[str, Any]]] = []

    def __init__(self, *args, **kwargs):
        pass

    async def __aenter__(self):
        return self

    async def __aexit__(self, exc_type, exc, tb):
        return False

    @classmethod
    def reset(cls, responses: list[httpx.Response]) -> None:
        cls.responses = list(responses)
        cls.calls = []

    @classmethod
    def _next(cls) -> httpx.Response:
        assert cls.responses, "unexpected HTTP request"
        return cls.responses.pop(0)

    async def post(self, url: str, **kwargs):
        self.calls.append(("POST", url, kwargs))
        return self._next()

    async def request(self, method: str, url: str, **kwargs):
        self.calls.append((method, url, kwargs))
        return self._next()

    class _Stream:
        def __init__(self, response: httpx.Response):
            self.response = response

        async def __aenter__(self):
            return self.response

        async def __aexit__(self, exc_type, exc, tb):
            return False

    def stream(self, method: str, url: str, **kwargs):
        self.calls.append(("DOWNLOAD", url, kwargs))
        return self._Stream(self._next())


def _response(
    *,
    json: Any = None,
    content: bytes | None = None,
    content_type: str = "application/json",
    url: str = "https://api.302.ai/test",
    status: int = 200,
) -> httpx.Response:
    request = httpx.Request("GET", url)
    if content is not None:
        return httpx.Response(status, content=content, headers={"content-type": content_type}, request=request)
    return httpx.Response(status, json=json, request=request)


def _result_zip(markdown: str) -> bytes:
    target = io.BytesIO()
    with zipfile.ZipFile(target, "w") as archive:
        archive.writestr("images/ignored.txt", "x")
        archive.writestr("full.md", markdown)
    return target.getvalue()


def _simple_pdf(text: str) -> bytes:
    """生成带可提取文本层的最小 PDF，避免测试依赖 PDF 写入库。"""
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, obj in enumerate(objects, 1):
        offsets.append(len(output))
        output.extend(f"{number} 0 obj\n".encode() + obj + b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode())
    output.extend(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode())
    return bytes(output)


def _simple_docx(path: Path, text: str) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
              <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
              <Default Extension="xml" ContentType="application/xml"/>
              <Override PartName="/word/document.xml"
                ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
            </Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1"
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"
                Target="word/document.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
            <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
              <w:body><w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:body>
            </w:document>""",
        )


def test_real_markitdown_converts_pdf_and_office_files(tmp_path):
    """依赖安装烟测：核心格式确实能产出可供引擎摄取的 Markdown。"""
    from openpyxl import Workbook
    from pptx import Presentation

    from sag_api.parsing.service import _markitdown_sync

    pdf = tmp_path / "sample.pdf"
    pdf.write_bytes(_simple_pdf("Muse PDF marker"))

    docx = tmp_path / "sample.docx"
    _simple_docx(docx, "Muse DOCX marker")

    pptx = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Muse PPTX marker"
    presentation.save(pptx)

    xlsx = tmp_path / "sample.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "Muse XLSX marker"
    workbook.save(xlsx)

    assert "Muse PDF marker" in _markitdown_sync(str(pdf))
    assert "Muse DOCX marker" in _markitdown_sync(str(docx))
    assert "Muse PPTX marker" in _markitdown_sync(str(pptx))
    assert "Muse XLSX marker" in _markitdown_sync(str(xlsx))


@pytest.mark.asyncio
async def test_mineru_upload_create_poll_and_download_zip(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset(
        [
            _response(json={"code": 200, "data": "https://file.302.ai/input.pdf", "message": "success"}),
            _response(json="task-123"),
            _response(json="processing"),
            _response(json='{"status":"done","full_zip_url":"https://file.302.ai/result.zip"}'),
            _response(
                content=_result_zip("# Parsed\n\nMinerU result"),
                content_type="application/zip",
                url="https://file.302.ai/result.zip",
            ),
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    states: list[dict[str, Any]] = []
    client = MinerUClient(
        _settings(
            mineru_base_url="https://api.302.ai",
            mineru_api_key="sk-mineru",
            mineru_poll_interval=0.001,
            mineru_poll_timeout=1,
        )
    )

    markdown = await client.parse(str(source), on_state=lambda state: _record(states, state))

    assert markdown == "# Parsed\n\nMinerU result\n"
    assert any(state.get("task_id") == "task-123" for state in states)
    assert [call[0] for call in _FakeAsyncClient.calls] == [
        "POST",
        "POST",
        "GET",
        "GET",
        "DOWNLOAD",
    ]
    assert _FakeAsyncClient.calls[0][1] == "https://api.302.ai/302/upload-file"
    assert _FakeAsyncClient.calls[1][2]["json"]["version"] == "2.5"
    assert "Authorization" not in _FakeAsyncClient.calls[-1][2]


@pytest.mark.asyncio
async def test_official_mineru_issue_105_upload_poll_and_download(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    upload_url = "https://upload.example.test/paper.pdf"
    result_url = "https://cdn-mineru.openxlab.org.cn/result.zip"
    _FakeAsyncClient.reset(
        [
            _response(
                json={
                    "code": 0,
                    "msg": "ok",
                    "trace_id": "trace-apply-105",
                    "data": {
                        "batch_id": "batch-105",
                        "file_urls": [upload_url],
                    },
                },
                url="https://mineru.net/api/v4/file-urls/batch",
            ),
            _response(content=b"", url=upload_url),
            _response(
                json={
                    "code": 0,
                    "msg": "ok",
                    "trace_id": "trace-poll-105",
                    "data": {
                        "batch_id": "batch-105",
                        "extract_result": [
                            {
                                "file_name": "paper.pdf",
                                "state": "done",
                                "err_msg": "",
                                "full_zip_url": result_url,
                            }
                        ],
                    },
                },
                url=(
                    "https://mineru.net/api/v4/extract-results/batch/batch-105"
                ),
            ),
            _response(
                content=_result_zip("# Official MinerU result"),
                content_type="application/zip",
                url=result_url,
            ),
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)

    async def allow_public_test_hosts(*_args, **_kwargs):
        return None

    monkeypatch.setattr(
        "sag_api.parsing.mineru._assert_public_host", allow_public_test_hosts
    )
    states: list[dict[str, Any]] = []
    official_settings = _settings(
        mineru_provider="official",
        mineru_base_url="https://mineru.net/api/v4/file-urls/batch",
        mineru_api_key="official-test-token",
        mineru_official_model="vlm",
        mineru_poll_interval=0.001,
        mineru_poll_timeout=1,
    )
    client = MinerUClient(official_settings)

    assert service._signature("mineru", official_settings) == (
        "mineru-official-vlm-auto"
    )

    markdown = await client.parse(
        str(source), on_state=lambda state: _record(states, state)
    )

    assert markdown == "# Official MinerU result\n"
    assert [call[:2] for call in _FakeAsyncClient.calls] == [
        ("POST", "https://mineru.net/api/v4/file-urls/batch"),
        ("PUT", upload_url),
        (
            "GET",
            "https://mineru.net/api/v4/extract-results/batch/batch-105",
        ),
        ("DOWNLOAD", result_url),
    ]
    assert _FakeAsyncClient.calls[0][2]["json"] == {
        "files": [{"name": "paper.pdf"}],
        "model_version": "vlm",
    }
    assert _FakeAsyncClient.calls[0][2]["headers"] == {
        "Authorization": "Bearer official-test-token"
    }
    assert "headers" not in _FakeAsyncClient.calls[1][2]
    assert "headers" not in _FakeAsyncClient.calls[3][2]
    assert any(
        state.get("batch_id") == "batch-105"
        and state.get("upload_completed") is True
        for state in states
    )


@pytest.mark.asyncio
async def test_official_mineru_resume_polls_existing_uploaded_batch(
    tmp_path, monkeypatch
):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    result_url = "https://cdn-mineru.openxlab.org.cn/resumed.zip"
    _FakeAsyncClient.reset(
        [
            _response(
                json={
                    "code": 0,
                    "msg": "ok",
                    "trace_id": "trace-resume-105",
                    "data": {
                        "batch_id": "batch-105",
                        "extract_result": [
                            {
                                "file_name": "paper.pdf",
                                "state": "done",
                                "err_msg": "",
                                "full_zip_url": result_url,
                            }
                        ],
                    },
                }
            ),
            _response(
                content=_result_zip("# Resumed official result"),
                content_type="application/zip",
                url=result_url,
            ),
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)

    async def allow_public_test_hosts(*_args, **_kwargs):
        return None

    monkeypatch.setattr(
        "sag_api.parsing.mineru._assert_public_host", allow_public_test_hosts
    )
    client = MinerUClient(
        _settings(
            mineru_provider="official",
            mineru_base_url="https://mineru.net",
            mineru_api_key="official-test-token",
            mineru_poll_interval=0.001,
            mineru_poll_timeout=1,
        )
    )

    markdown = await client.parse(
        str(source),
        state={
            "mineru_service": "official",
            "batch_id": "batch-105",
            "upload_completed": True,
            "filename": "paper.pdf",
        },
    )

    assert markdown == "# Resumed official result\n"
    assert [call[:2] for call in _FakeAsyncClient.calls] == [
        (
            "GET",
            "https://mineru.net/api/v4/extract-results/batch/batch-105",
        ),
        ("DOWNLOAD", result_url),
    ]


@pytest.mark.asyncio
async def test_official_mineru_business_error_keeps_message_and_trace(
    tmp_path, monkeypatch
):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset(
        [
            _response(
                json={
                    "code": -60011,
                    "msg": "请确保文件已上传",
                    "trace_id": "trace-105",
                    "data": None,
                }
            )
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)

    async def allow_public_test_hosts(*_args, **_kwargs):
        return None

    monkeypatch.setattr(
        "sag_api.parsing.mineru._assert_public_host", allow_public_test_hosts
    )
    client = MinerUClient(
        _settings(
            mineru_provider="official",
            mineru_base_url="https://mineru.net",
            mineru_api_key="official-test-token",
        )
    )

    with pytest.raises(
        UpstreamError, match="请确保文件已上传.*trace_id: trace-105"
    ):
        await client.parse(str(source))


async def _record(target: list[dict[str, Any]], state: dict[str, Any]) -> None:
    target.append(dict(state))


@pytest.mark.asyncio
async def test_mineru_retry_reuses_upload_and_task_id(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset([_response(json="# Ready\n\nNo duplicate paid task")])
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    client = MinerUClient(
        _settings(
            mineru_base_url="https://api.302.ai",
            mineru_api_key="sk-mineru",
            mineru_poll_interval=0.001,
            mineru_poll_timeout=1,
        )
    )

    markdown = await client.parse(
        str(source),
        state={"upload_url": "https://file.302.ai/input.pdf", "task_id": "existing-task"},
    )
    assert markdown.startswith("# Ready")
    assert [call[0] for call in _FakeAsyncClient.calls] == ["GET"]


@pytest.mark.asyncio
async def test_mineru_accepts_immediate_result_url_from_create(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset(
        [
            _response(json={"code": 200, "data": "https://file.302.ai/input.pdf"}),
            _response(json="https://file.302.ai/result.zip"),
            _response(
                content=_result_zip("# Immediate result"),
                content_type="application/zip",
                url="https://file.302.ai/result.zip",
            ),
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    client = MinerUClient(_settings(mineru_base_url="https://api.302.ai", mineru_api_key="sk-mineru"))

    assert await client.parse(str(source)) == "# Immediate result\n"
    assert [call[0] for call in _FakeAsyncClient.calls] == ["POST", "POST", "DOWNLOAD"]


@pytest.mark.asyncio
async def test_mineru_result_download_maps_bounded_error_body(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    oversized_error = {"data": {"err_msg": "result temporarily unavailable"}}
    _FakeAsyncClient.reset(
        [
            _response(json={"code": 200, "data": "https://file.302.ai/input.pdf"}),
            _response(json="https://file.302.ai/result.md"),
            _response(json=oversized_error, status=503),
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    client = MinerUClient(_settings(mineru_base_url="https://api.302.ai", mineru_api_key="sk-mineru"))

    with pytest.raises(ServiceUnavailableError, match="result temporarily unavailable"):
        await client.parse(str(source))


@pytest.mark.asyncio
async def test_mineru_auth_error_is_configuration_error(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset([_response(json={"message": "bad key"}, status=401)])
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    client = MinerUClient(_settings(mineru_base_url="https://api.302.ai", mineru_api_key="sk-bad"))
    with pytest.raises(ConfigurationError, match="API Key"):
        await client.parse(str(source))


def test_mineru_zip_requires_markdown():
    target = io.BytesIO()
    with zipfile.ZipFile(target, "w") as archive:
        archive.writestr("result.json", "{}")
    from sag_api.parsing.mineru import _markdown_from_zip

    with pytest.raises(UpstreamError, match="没有 Markdown"):
        _markdown_from_zip(target.getvalue(), 1024)


def test_mineru_pending_and_nested_failure_payloads_are_not_misclassified():
    pending = {
        "status": "processing",
        "data": {"url": "https://file.302.ai/input.pdf"},
    }
    assert _interpret_poll_payload(pending, "task-1")[0] == "pending"

    nested_failure = {
        "code": 200,
        "msg": "ok",
        "data": {"status": "failed", "err_msg": "bad pdf"},
    }
    assert _interpret_poll_payload(nested_failure, "task-1") == ("failed", "bad pdf")
    assert _interpret_poll_payload("A0202", "task-1")[0] == "failed"


@pytest.mark.asyncio
async def test_mineru_result_download_rejects_private_hosts():
    with pytest.raises(UpstreamError, match="内网"):
        await _assert_public_host("127.0.0.1", 80)
    with pytest.raises(UpstreamError, match="内网"):
        await _assert_public_host("169.254.169.254", 80)


@pytest.mark.asyncio
@pytest.mark.parametrize("fake_ip", ["198.18.0.90", "198.19.255.254"])
async def test_mineru_result_download_allows_fake_ip_proxy_dns(fake_ip):
    # Clash/mihomo 等代理的 fake-ip 模式会把域名解析到 198.18.0.0/15，
    # TUN 网关按映射表转发到真实公网地址；该段不指向任何真实内网服务，
    # 不应被 SSRF 守卫当作内网拒绝。
    await _assert_public_host(fake_ip, 443)
